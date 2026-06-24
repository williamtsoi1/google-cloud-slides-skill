#!/usr/bin/env python3
"""build_gslides.py — Build a Google Cloud-branded .pptx from a JSON deck spec.

The output .pptx is designed to be uploaded to Google Drive and opened as Google
Slides, where Drive converts it to a native, editable deck (File -> Save as Google
Slides, or with Drive's "convert uploads" setting on). There is intentionally NO
network/auth/upload here — the upload is a manual step (see SKILL.md).

Brand notes:
  * Font is Roboto, NOT Google Sans. Google Sans is proprietary and is not in the
    Google Slides font library, and Slides ignores fonts embedded in a .pptx on
    import — so an editable deck cannot render real Google Sans. Roboto is Google's
    own typeface, is always available in Slides, and is the closest free relative.
    (The HTML/PDF export path keeps real Google Sans.)
  * The cloud logo and rainbow bar are the bundled image assets, embedded directly
    from templates/ — never recreated with shapes/gradients.
  * Every slide gets the footer: "Google Cloud" bottom-left, the confidential
    notice + page number bottom-right. Section numbers are zero-padded (01, 02).

Usage:
  build_gslides.py <deck.json> [out.pptx]

See references/GSLIDES_SPEC.md for the deck.json schema and supported slide types.
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:  # pragma: no cover - surfaced via the .sh wrapper
    sys.exit(
        "build_gslides: python-pptx is not installed.\n"
        "  Run via build_gslides.sh (auto-installs with uv), or: pip install python-pptx"
    )

# --- Brand constants ---------------------------------------------------------
FONT = "Roboto"  # see module docstring for why this is not "Google Sans"

BLUE = RGBColor(0x31, 0x86, 0xFF)
RED = RGBColor(0xFC, 0x41, 0x3D)
YELLOW = RGBColor(0xFE, 0xC7, 0x00)
GREEN = RGBColor(0x00, 0xAF, 0x57)
DARK = RGBColor(0x20, 0x21, 0x24)
LIGHT = RGBColor(0xFF, 0xFF, 0xFF)
FOOTER_GRAY = RGBColor(0x5F, 0x63, 0x68)

SECTION_BG = {"green": GREEN, "blue": BLUE, "red": RED}

SCRIPT_DIR = Path(__file__).resolve().parent
TEMPLATES = SCRIPT_DIR.parent / "templates"
CLOUD_LOGO = TEMPLATES / "gradient_super_cloud_512_2x.png"
RAINBOW_BAR = TEMPLATES / "GC_Progress_Bar_Gradient_RGB.jpg"

# 16:9 at the python-pptx widescreen size.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Resolved at load time so slide-referenced images can be relative to the deck file.
DECK_DIR = Path(".")


# --- Low-level helpers -------------------------------------------------------
def _norm_runs(runs):
    """Normalize a paragraph's runs into a list of dicts with a 'text' key."""
    if isinstance(runs, str):
        return [{"text": runs}]
    out = []
    for r in runs:
        out.append({"text": r} if isinstance(r, str) else dict(r))
    return out


def add_text(
    slide,
    left,
    top,
    width,
    height,
    paragraphs,
    *,
    size=18,
    bold=False,
    color=DARK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=None,
    space_after=None,
):
    """Add a textbox. `paragraphs` is a single paragraph (str / list of runs) or a
    list of paragraphs. Each run may be a str or a dict overriding text/color/size/bold."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    # Detect: list of paragraphs vs a single paragraph expressed as a list of runs.
    if isinstance(paragraphs, str):
        para_list = [paragraphs]
    elif paragraphs and isinstance(paragraphs[0], (str, dict)):
        para_list = [paragraphs]  # one paragraph, multiple runs
    else:
        para_list = list(paragraphs)

    for i, para in enumerate(para_list):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        if space_after is not None:
            p.space_after = Pt(space_after)
        for spec in _norm_runs(para):
            run = p.add_run()
            run.text = spec["text"]
            f = run.font
            f.name = FONT
            f.size = Pt(spec.get("size", size))
            f.bold = spec.get("bold", bold)
            f.color.rgb = spec.get("color", color)
    return tb


def add_fullbleed(slide, color):
    """Full-slide solid background rectangle (added first so text sits on top)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_rainbow(slide, left, top, width):
    """The brand rainbow divider bar, from the bundled image asset."""
    slide.shapes.add_picture(str(RAINBOW_BAR), left, top, width=width, height=Inches(0.06))


def add_cloud(slide, left, top, size=Inches(4.0)):
    """The official Google Cloud icon, from the bundled image asset (square)."""
    slide.shapes.add_picture(str(CLOUD_LOGO), left, top, width=size, height=size)


def resolve_image(path_str):
    p = Path(path_str)
    if not p.is_absolute():
        p = DECK_DIR / p
    if not p.exists():
        raise FileNotFoundError(f"image not found: {path_str} (resolved to {p})")
    return str(p)


def add_footer(slide, page_no, footer_text, *, dark_bg=False):
    """'Google Cloud' bottom-left, '<notice>  <page>' bottom-right. On every slide."""
    color = LIGHT if dark_bg else FOOTER_GRAY
    add_text(
        slide, Inches(0.55), Inches(7.0), Inches(4), Inches(0.35),
        "Google Cloud", size=10, color=color,
    )
    add_text(
        slide, Inches(8.78), Inches(7.0), Inches(4), Inches(0.35),
        f"{footer_text}  {page_no}", size=10, color=color, align=PP_ALIGN.RIGHT,
    )


# --- Layout renderers --------------------------------------------------------
# Each takes (slide, spec) and draws the body; the caller adds the footer.

def render_cover(slide, s):
    add_text(slide, Inches(0.7), Inches(0.5), Inches(6), Inches(0.4),
             "Google Cloud", size=14, bold=True, color=DARK)
    add_cloud(slide, Inches(8.6), Inches(1.85))
    add_text(slide, Inches(0.7), Inches(2.3), Inches(7.5), Inches(2.6),
             s["title"], size=44, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.05)
    sub = s.get("subtitle")
    date = s.get("date")
    extra = "  •  ".join(x for x in (sub, date) if x)
    if extra:
        add_text(slide, Inches(0.7), Inches(5.0), Inches(7.5), Inches(0.6),
                 extra, size=18, color=BLUE)


def render_section(slide, s):
    number = f"{int(s['number']):02d}"
    color = (s.get("color") or "white").lower()
    if color == "white":
        add_rainbow(slide, Inches(0.7), Inches(0.95), Inches(11.9))
        add_text(slide, Inches(0.7), Inches(2.4), Inches(12), Inches(1.6),
                 number, size=96, bold=True, color=DARK)
        add_text(slide, Inches(0.7), Inches(4.2), Inches(12), Inches(1.2),
                 s["title"], size=36, bold=True, color=DARK)
        return False  # light bg
    add_fullbleed(slide, SECTION_BG.get(color, GREEN))
    add_text(slide, Inches(0.7), Inches(0.7), Inches(4), Inches(1.4),
             number, size=80, bold=True, color=LIGHT)
    add_text(slide, Inches(0.7), Inches(2.6), Inches(11.9), Inches(1.4),
             s["title"], size=40, bold=True, color=LIGHT, anchor=MSO_ANCHOR.MIDDLE)
    items = s.get("items") or []
    if items:
        add_text(slide, Inches(0.7), Inches(4.2), Inches(11.9), Inches(2.2),
                 [[{"text": it}] for it in items], size=18, color=LIGHT,
                 line_spacing=1.2, space_after=8)
    return True  # dark/colored bg


def _title(slide, text):
    add_text(slide, Inches(0.7), Inches(0.6), Inches(11.9), Inches(1.0),
             text, size=30, bold=True, color=DARK)


def render_title_body(slide, s):
    _title(slide, s["title"])
    add_text(slide, Inches(0.7), Inches(1.9), Inches(11.9), Inches(4.6),
             s.get("body", ""), size=18, color=DARK, line_spacing=1.3)


def render_two_column(slide, s):
    _title(slide, s["title"])
    cols = s["columns"]
    add_text(slide, Inches(0.7), Inches(1.9), Inches(5.7), Inches(4.6),
             cols[0], size=18, color=DARK, line_spacing=1.3)
    if len(cols) > 1:
        add_text(slide, Inches(6.9), Inches(1.9), Inches(5.7), Inches(4.6),
                 cols[1], size=18, color=DARK, line_spacing=1.3)


def render_bullets(slide, s):
    _title(slide, s["title"])
    paras = [[{"text": f"•   {b}"}] for b in s["bullets"]]
    add_text(slide, Inches(0.9), Inches(1.9), Inches(11.5), Inches(4.6),
             paras, size=22, color=DARK, line_spacing=1.2, space_after=12)


def render_stat(slide, s):
    add_text(slide, Inches(0.7), Inches(1.8), Inches(11.9), Inches(2.0),
             s["stat"], size=120, bold=True, color=BLUE)
    add_text(slide, Inches(0.75), Inches(4.0), Inches(11.9), Inches(0.8),
             s.get("label", ""), size=28, bold=True, color=DARK)
    if s.get("support"):
        add_text(slide, Inches(0.75), Inches(4.9), Inches(11.9), Inches(1.2),
                 s["support"], size=18, color=DARK, line_spacing=1.3)


def render_three_stat(slide, s):
    _title(slide, s["title"])
    stats = s["stats"][:3]
    width = Inches(3.9)
    for i, st in enumerate(stats):
        left = Inches(0.7 + i * 4.1)
        add_text(slide, left, Inches(2.4), width, Inches(1.4),
                 st["stat"], size=72, bold=True, color=BLUE)
        add_text(slide, left, Inches(4.0), width, Inches(1.4),
                 st.get("label", ""), size=18, color=DARK, line_spacing=1.2)


def render_statement(slide, s):
    dark = bool(s.get("dark"))
    if dark:
        add_fullbleed(slide, DARK)
    base = LIGHT if dark else DARK
    text = s["text"]
    highlight = s.get("highlight")
    if highlight and text.startswith(highlight):
        runs = [{"text": highlight, "color": BLUE},
                {"text": text[len(highlight):], "color": base}]
    elif highlight:
        runs = [{"text": highlight + " ", "color": BLUE}, {"text": text, "color": base}]
    else:
        runs = [{"text": text, "color": base}]
    add_text(slide, Inches(0.9), Inches(1.2), Inches(11.5), Inches(5.0),
             runs, size=48, bold=False, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    if dark:
        add_rainbow(slide, Inches(0.9), Inches(6.3), Inches(4.0))
    return dark


def render_two_tone(slide, s):
    dark = bool(s.get("dark"))
    if dark:
        add_fullbleed(slide, DARK)
    base = LIGHT if dark else DARK
    add_text(slide, Inches(0.9), Inches(1.6), Inches(11.5), Inches(4.4),
             [[{"text": s["line1"], "color": base}],
              [{"text": s["line2"], "color": BLUE}]],
             size=60, bold=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    return dark


def render_quote(slide, s):
    add_fullbleed(slide, DARK)
    add_text(slide, Inches(0.9), Inches(0.7), Inches(11.5), Inches(1.4),
             "“", size=120, bold=True, color=BLUE)
    add_text(slide, Inches(0.9), Inches(2.2), Inches(11.5), Inches(3.0),
             s["quote"], size=36, color=LIGHT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    if s.get("attribution"):
        add_text(slide, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.8),
                 s["attribution"], size=20, bold=True, color=BLUE)
    return True


def render_image_headline(slide, s):
    side = (s.get("side") or "left").lower()
    img = resolve_image(s["image"])
    half = Inches(6.6667)
    if side == "left":
        slide.shapes.add_picture(img, 0, 0, width=half, height=SLIDE_H)
        tx_left = Inches(7.2)
    else:
        slide.shapes.add_picture(img, half, 0, width=half, height=SLIDE_H)
        tx_left = Inches(0.7)
    add_text(slide, tx_left, Inches(1.0), Inches(5.4), Inches(5.0),
             s["headline"], size=40, bold=True, color=DARK,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)


def render_bullets_image(slide, s):
    _title(slide, s["title"])
    paras = [[{"text": f"•   {b}"}] for b in s["bullets"]]
    add_text(slide, Inches(0.9), Inches(1.9), Inches(6.6), Inches(4.6),
             paras, size=20, color=DARK, line_spacing=1.2, space_after=10)
    img = resolve_image(s["image"])
    slide.shapes.add_picture(img, Inches(8.0), Inches(1.9), width=Inches(4.6), height=Inches(4.4))


def render_thank_you(slide, s):
    add_cloud(slide, Inches(8.6), Inches(1.85))
    add_text(slide, Inches(0.7), Inches(2.3), Inches(7.5), Inches(2.0),
             s.get("title", "Thank you"), size=64, bold=True, color=DARK,
             anchor=MSO_ANCHOR.MIDDLE)


RENDERERS = {
    "cover": render_cover,
    "section": render_section,
    "title-body": render_title_body,
    "two-column": render_two_column,
    "bullets": render_bullets,
    "stat": render_stat,
    "three-stat": render_three_stat,
    "statement": render_statement,
    "two-tone": render_two_tone,
    "quote": render_quote,
    "image-headline": render_image_headline,
    "bullets-image": render_bullets_image,
    "thank-you": render_thank_you,
}


# --- Driver ------------------------------------------------------------------
def build(deck, out_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    footer_text = deck.get("footer", "Proprietary & Confidential")
    slides = deck.get("slides", [])
    if not slides:
        raise ValueError("deck has no slides")

    for i, spec in enumerate(slides, start=1):
        stype = spec.get("type")
        renderer = RENDERERS.get(stype)
        if renderer is None:
            supported = ", ".join(sorted(RENDERERS))
            raise ValueError(
                f"slide {i}: unsupported type {stype!r}. Supported types: {supported}"
            )
        slide = prs.slides.add_slide(blank)
        try:
            dark_bg = bool(renderer(slide, spec))
        except KeyError as e:
            raise ValueError(
                f"slide {i} (type {stype!r}): missing required field {e}"
            ) from None
        add_footer(slide, i, footer_text, dark_bg=dark_bg)

    prs.save(out_path)
    return len(slides)


def main(argv):
    global DECK_DIR
    if len(argv) < 2:
        sys.exit("usage: build_gslides.py <deck.json> [out.pptx]")
    deck_path = Path(argv[1])
    if not deck_path.exists():
        sys.exit(f"build_gslides: deck spec not found: {deck_path}")
    DECK_DIR = deck_path.resolve().parent

    out_path = Path(argv[2]) if len(argv) >= 3 else deck_path.with_suffix(".pptx")

    try:
        deck = json.loads(deck_path.read_text())
    except json.JSONDecodeError as e:
        sys.exit(f"build_gslides: invalid JSON in {deck_path}: {e}")

    try:
        n = build(deck, str(out_path))
    except (ValueError, FileNotFoundError) as e:
        sys.exit(f"build_gslides: {e}")

    print(f"build_gslides: wrote {out_path} ({n} slides)", file=sys.stderr)
    print(
        "Next: upload it to Google Drive and open it as Google Slides "
        "(File -> Save as Google Slides) to get a native editable deck.",
        file=sys.stderr,
    )


if __name__ == "__main__":
    main(sys.argv)
